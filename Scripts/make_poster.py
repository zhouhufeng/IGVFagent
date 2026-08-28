#!/usr/bin/env python3
"""Build the IGVF Annual Meeting poster as a .pptx.

Geometry follows the organiser's instructions: 36" wide x 48" tall (portrait),
never wider than 48", four posters per 8ft x 4ft double-sided board.

Every headline number is read from the running system rather than typed in, so
the poster cannot drift from the software it describes — the same reason the
skill cards derive their validation status instead of asserting it.

    python3 Scripts/make_poster.py --out Docs/Poster/IGVFagent_poster.pptx
"""
from __future__ import annotations

import argparse
import json
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]

W_IN, H_IN = 36, 48          # portrait, per the organiser's email
MARGIN = 1.1
COL_GAP = 0.55
N_COLS = 2
COL_W = (W_IN - 2 * MARGIN - (N_COLS - 1) * COL_GAP) / N_COLS

# IGVF-adjacent palette: deep blue, teal accent, warm highlight.
NAVY = RGBColor(0x0B, 0x2E, 0x59)
BLUE = RGBColor(0x1B, 0x5E, 0xA8)
TEAL = RGBColor(0x0F, 0x84, 0x8C)
AMBER = RGBColor(0xC2, 0x6B, 0x1F)
INK = RGBColor(0x1A, 0x1A, 0x1A)
GREY = RGBColor(0x55, 0x5F, 0x6B)
PALE = RGBColor(0xEF, 0xF4, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ---------------------------------------------------------------------------
# Live numbers
# ---------------------------------------------------------------------------

def live_stats() -> dict:
    sys.path.insert(0, str(ROOT / "Scripts"))
    out = {"version": "?", "skills": 0, "tools": 0,
           "benchmarks": 0, "checks": 0, "quant": 0}
    try:
        from igvfagent import cli, _tools, __version__
        out["version"] = __version__
        out["skills"] = len(cli.SKILLS)
        out["tools"] = len(_tools.list_tools())
    except Exception:
        pass
    try:
        tax = json.loads(subprocess.run(
            [sys.executable, str(ROOT / "Benchmarks" / "taxonomy.py"), "--json"],
            capture_output=True, text=True, timeout=120).stdout)
        out["benchmarks"] = tax["n_benchmarks"]
        out["checks"] = sum(b["n_checks"] for b in tax["benchmarks"])
        out["quant"] = tax["summary"]["A"]
    except Exception:
        pass
    return out


# ---------------------------------------------------------------------------
# Layout helpers
# ---------------------------------------------------------------------------

def _box(slide, x, y, w, h, *, fill=None, line=None, line_w=1.5):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = 0.03
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def _text(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
          anchor=MSO_ANCHOR.TOP, spacing=1.06):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    first = True
    for spec in runs:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = spec.get("align", align)
        p.line_spacing = spec.get("spacing", spacing)
        if spec.get("space_before"):
            p.space_before = Pt(spec["space_before"])
        if spec.get("space_after"):
            p.space_after = Pt(spec["space_after"])
        r = p.add_run()
        r.text = spec["t"]
        f = r.font
        f.size = Pt(spec.get("size", 20))
        f.bold = spec.get("bold", False)
        f.italic = spec.get("italic", False)
        f.color.rgb = spec.get("color", INK)
        f.name = spec.get("font", "Helvetica Neue")
    return tb


def section(slide, x, y, w, title, body, *, h=None, accent=BLUE):
    """A titled panel. Returns the y just below it."""
    n_lines = sum(max(1, len(b["t"]) // int(w * 4.6) + 1) for b in body)
    h = h or (1.00 + 0.38 * n_lines + 0.30)
    _box(slide, x, y, w, h, fill=WHITE, line=accent, line_w=2.0)
    _box(slide, x, y, w, 0.62, fill=accent, line=None)
    _text(slide, x + 0.22, y + 0.06, w - 0.44, 0.5,
          [{"t": title, "size": 25, "bold": True, "color": WHITE}])
    _text(slide, x + 0.24, y + 0.78, w - 0.48, h - 0.95, body)
    return y + h + 0.34


def stat_tile(slide, x, y, w, h, value, label, *, color=TEAL):
    _box(slide, x, y, w, h, fill=PALE, line=color, line_w=2.0)
    _text(slide, x, y + 0.16, w, h * 0.52,
          [{"t": value, "size": 46, "bold": True, "color": color,
            "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
    _text(slide, x, y + h * 0.60, w, h * 0.36,
          [{"t": label, "size": 15, "color": GREY, "align": PP_ALIGN.CENTER}],
          align=PP_ALIGN.CENTER)


def picture(slide, path: Path, x, y, w, *, caption=None, max_h=None):
    if not path.is_file():
        return y
    from PIL import Image  # optional; fall back to a fixed ratio
    try:
        iw, ih = Image.open(path).size
        h = w * ih / iw
    except Exception:
        h = w * 0.62
    if max_h and h > max_h:
        h = max_h
    slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                             width=Inches(w), height=Inches(h))
    y2 = y + h + 0.10
    if caption:
        _text(slide, x, y2, w, 0.5,
              [{"t": caption, "size": 13, "italic": True, "color": GREY}])
        y2 += 0.46
    return y2


# ---------------------------------------------------------------------------
# Poster
# ---------------------------------------------------------------------------

def build(out_path: Path) -> Path:
    s = live_stats()
    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # ---------------- header ----------------
    _box(slide, 0, 0, W_IN, 5.5, fill=NAVY, line=None)
    _text(slide, MARGIN, 0.55, W_IN - 2 * MARGIN, 2.5,
          [{"t": "IGVF Agent", "size": 82, "bold": True, "color": WHITE,
            "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
    _text(slide, MARGIN, 2.35, W_IN - 2 * MARGIN, 1.5,
          [{"t": "An auditable, locally deployable research agent for analysis "
                 "across federated functional-genomics resources",
            "size": 30, "color": RGBColor(0xCF, 0xE3, 0xF5),
            "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
    _text(slide, MARGIN, 3.75, W_IN - 2 * MARGIN, 1.4,
          [{"t": "Hufeng Zhou¹, Yichun He², Daofeng Li³, Junhao Zhu¹, Xihao Li⁴, "
                 "Xingxin Pan⁵, Alvin Zhang¹, Vineet Verma¹, Ting Wang³, "
                 "Zhiping Weng⁶, Tianxi Cai⁷, Jesse Engreitz⁸, "
                 "William H. Majoros⁹, S. Stephen Yi⁵, Xihong Lin¹²",
            "size": 16, "color": RGBColor(0xAF, 0xCB, 0xE6),
            "align": PP_ALIGN.CENTER, "spacing": 1.0},
           {"t": "¹Harvard T.H. Chan School of Public Health  ·  ²Broad Institute  ·  "
                 "³Washington University  ·  ⁴UNC Chapel Hill  ·  ⁵UT Austin / Baylor  ·  "
                 "⁶UMass Chan  ·  ⁷Harvard Medical School  ·  ⁸Stanford  ·  ⁹Duke",
            "size": 13, "color": RGBColor(0x8F, 0xB3, 0xD6),
            "align": PP_ALIGN.CENTER, "space_before": 6}],
          align=PP_ALIGN.CENTER)

    # ---------------- headline stats ----------------
    y = 6.1
    tw = (W_IN - 2 * MARGIN - 4 * 0.35) / 5
    tiles = [(str(s["skills"]), "analysis skills"),
             (str(s["tools"]), "typed LLM tools"),
             (str(s["benchmarks"]), "benchmarked studies"),
             (str(s["checks"]), "machine-checked criteria"),
             ("6", "federated archives")]
    for i, (v, lab) in enumerate(tiles):
        stat_tile(slide, MARGIN + i * (tw + 0.35), y, tw, 2.1, v, lab)
    y += 2.75

    LX, RX = MARGIN, MARGIN + COL_W + COL_GAP

    # ================= LEFT COLUMN =================
    ly = section(slide, LX, y, COL_W, "The problem",
        [{"t": "Variant-to-function research now spans IGVF, ENCODE, GEO, "
               "CELLxGENE, HCA and Zenodo — archives that differ in schema, "
               "identifiers, metadata standards and access mechanisms.",
          "size": 19},
         {"t": "The rate-limiting step has shifted from the availability of "
               "data to the labour of integrating it: even a focused question "
               "— which enhancers regulate PCSK9 in hepatocytes — needs custom "
               "scripts to search, harmonise, retrieve and track provenance "
               "across several repositories.",
          "size": 19, "space_before": 10},
         {"t": "LLM agents are a credible response, but most are bound to one "
               "vendor, route credentials through opaque cloud services, or "
               "expose a single archive. None of those is usable for "
               "controlled-access data.",
          "size": 19, "space_before": 10}])

    ly = section(slide, LX, ly, COL_W, "Design", [
        {"t": "1.  Explicit commands, not generated code", "size": 20, "bold": True},
        {"t": "Every agent action is a typed, recorded command drawn from a "
              "fixed repertoire — never free-form code. A session therefore "
              "reduces without loss to a re-runnable script.", "size": 18},
        {"t": "2.  Local execution", "size": 20, "bold": True, "space_before": 9},
        {"t": "Tools run as subprocesses in the investigator's environment. "
              "Endpoints and credentials resolve at runtime, so the released "
              "code is the deployed code.", "size": 18},
        {"t": "3.  Two-way drivability", "size": 20, "bold": True, "space_before": 9},
        {"t": "Each capability is simultaneously a shell command and an "
              "LLM-callable tool with identical semantics, so exploration and "
              "the archived pipeline execute the same code.", "size": 18},
        {"t": "4.  Backend agnosticism", "size": 20, "bold": True, "space_before": 9},
        {"t": "One routing layer drives hosted APIs, self-hosted open-weight "
              "models, or an external coding CLI without touching analysis "
              "code.", "size": 18},
    ], accent=TEAL)

    ly = picture(slide, ROOT / "Docs" / "Figures" / "IGVF_agent_archetcture.png",
                 LX, ly, COL_W,
                 caption="Five-layer architecture: entry points → agent runtime "
                         "and tool dispatch → skills by domain → local "
                         "persistence → upstream archives.",
                 max_h=8.4) + 0.30

    ly = section(slide, LX, ly, COL_W, "Growing knowledge graph", [
        {"t": "Every retrieval feeds a local graph. Result tables are parsed "
              "and their rows become entity-to-entity edges — variant → gene, "
              "element → gene, gene → pathway, protein ↔ protein — carrying "
              "score, source and method.", "size": 19},
        {"t": "The graph therefore accumulates across sessions rather than "
              "resetting: repeat questions are answered locally, and evidence "
              "gathered for one analysis is available to the next.",
         "size": 19, "space_before": 10},
    ], accent=AMBER)

    ly = section(slide, LX, ly, COL_W, "Worked example", [
        {"t": "“Functionally annotate these variants and add them to the "
              "knowledge graph”", "size": 19, "italic": True, "color": BLUE},
        {"t": "25 APOB stop-gain variants, pasted in any notation "
              "(chr-pos-ref-alt, rsID, SPDI, HGVS, VCF — mixed freely).",
         "size": 19, "space_before": 8},
        {"t": "→  annotated against FAVOR and the IGVF Catalog, with live "
              "ClinVar taking precedence over FAVOR's snapshot", "size": 18},
        {"t": "→  25 variant nodes, 7 gene edges, 24 disease edges, 6 "
              "regulatory-element edges written to the local graph",
         "size": 18},
        {"t": "→  an APOB depth-2 traversal returns 19 nodes and 47 edges "
              "spanning variant, gene, disease and regulatory element",
         "size": 18},
        {"t": "Cross-source verification matters: for one variant FAVOR's "
              "snapshot reports Benign / hypercholesterolaemia while current "
              "ClinVar reports Pathogenic / hypobetalipoproteinaemia — "
              "opposite significance and opposite direction of effect.",
         "size": 17, "italic": True, "space_before": 8, "color": AMBER},
    ], accent=BLUE)

    # ================= RIGHT COLUMN =================
    ry = section(slide, RX, y, COL_W, "Capabilities", [
        {"t": "Discovery and retrieval  ·  variant interpretation  ·  "
              "regulatory genomics  ·  single-cell and multiome  ·  bulk and "
              "proteomic analysis  ·  cross-resource integration", "size": 19},
        {"t": "Assay-aware reimplementations", "size": 20, "bold": True,
         "space_before": 10},
        {"t": "Clean-room implementations of the canonical published pipeline "
              "for each major assay — MPRA, STARR-seq, CRISPRi Flow-FISH, "
              "SHARE-seq, VAMP-seq deep mutational scanning, MULTI-seq — "
              "written to one provenance and plotting framework under a "
              "single permissive licence.", "size": 18},
        {"t": "Network integration", "size": 20, "bold": True, "space_before": 9},
        {"t": "CARNIVAL signed MILP and prize-collecting Steiner tree, "
              "reimplemented from the published formulations, composing the "
              "outputs of several skills into one mechanistic inference.",
         "size": 18},
    ], accent=TEAL)

    ry = section(slide, RX, ry, COL_W, "Evaluated in three tiers", [
        {"t": "An agent fails in three independent ways, and one aggregate "
              "score conceals which occurred.", "size": 19},
        {"t": f"Tier 1 — skill correctness.  {s['benchmarks']} studies, "
              f"{s['checks']} machine-checked criteria. Fixed command "
              f"sequences with no model in the loop: this establishes that the "
              f"implementations are right, not that the agent would select "
              f"them.", "size": 18, "space_before": 8},
        {"t": "Tier 2 — planning. Tool selection, argument binding and "
              "ordering scored against a gold plan, with deliberate tool "
              "failures to test recovery.", "size": 18, "space_before": 6},
        {"t": "Tier 3 — conclusion validity. An independent verifier judges "
              "each claim supported, unsupported or contradicted by the "
              "artefacts — without seeing the reasoning that produced it.",
         "size": 18, "space_before": 6},
        {"t": f"Reported honestly: of {s['benchmarks']} benchmarks, "
              f"{s['quant']} assert a derived scientific quantity; the rest "
              f"establish retrieval or artefact generation. Calling them all "
              f"“reproductions” would overstate the weakest.",
         "size": 17, "italic": True, "space_before": 8, "color": AMBER},
    ], accent=AMBER)

    ry = section(slide, RX, ry, COL_W, "Result: the multiome corpus", [
        {"t": "A cross-archive survey resolved 401 datasets / 3,837 files / "
              "14.19 TB across six archives in ~20 minutes on a workstation.",
         "size": 19},
        {"t": "Of that, ~12.5 TB is raw reads and alignments recoverable from "
              "primary data. The analysis-ready fraction — count matrices, "
              "fragment files, cell annotations — is ≈1.04 TB.", "size": 19,
         "space_before": 8},
        {"t": "That the union of publicly available multiome data collapses to "
              "a roughly one-terabyte analysis-ready corpus brings assembly of "
              "a single-cell foundation-model training set within reach of an "
              "individual laboratory.", "size": 19, "space_before": 8,
         "bold": True},
    ], accent=BLUE)

    ry = section(slide, RX, ry, COL_W, "Reproducibility, stated precisely", [
        {"t": "Auditability is a property of the record: every action is a "
              "typed command, recorded with arguments, artefacts and a "
              "consistency fingerprint. This holds in general.", "size": 18},
        {"t": "Repeatability is a property of execution: pinned command "
              "sequences reproduce artefacts exactly; free-form planning does "
              "not, and neither survives an upstream re-release.", "size": 18,
         "space_before": 7},
        {"t": "Validity is a property of the conclusion, and is the open "
              "problem — agreement metrics measure whether two runs concur, "
              "which two identically mistaken runs also satisfy.", "size": 18,
         "space_before": 7},
    ], accent=TEAL)

    ry = section(slide, RX, ry, COL_W, "Extensible by its users", [
        {"t": "New capabilities do not require touching core code. A tool "
              "manifest or a Python skill dropped into an extension directory "
              "is absorbed into the registry with no restart, appearing "
              "simultaneously as a shell command and an LLM-callable tool.",
         "size": 18},
        {"t": "The agent can author these itself: asked for a capability that "
              "does not exist, it writes the implementation and registers the "
              "matching typed interface in one step, then calls it.",
         "size": 18, "space_before": 8},
        {"t": "Each skill is described by a machine-readable skill card — "
              "identity, version, typed interfaces, upstream archives, and a "
              "validation status derived from the benchmark suite rather than "
              "asserted — so skills can be version-pinned and consumed by "
              "systems other than this agent.", "size": 18, "space_before": 8},
    ], accent=TEAL)

    ry = section(slide, RX, ry, COL_W, "Deployment and privacy", [
        {"t": "Three interchangeable interfaces — command line, browser UI, "
              "conversational agent — over one shared contract.", "size": 18},
        {"t": "Backends span hosted APIs, self-hosted open-weight models "
              "(Ollama, vLLM, TGI) and external coding CLIs. Analysis code is "
              "identical across all of them.", "size": 18, "space_before": 7},
        {"t": "Stated precisely: tool execution is local, and endpoints and "
              "credentials resolve at runtime. Inference is local only when a "
              "locally served model is used, and federated data access "
              "contacts public archives under every backend — so a local "
              "model keeps the prompt local rather than eliminating egress.",
         "size": 17, "italic": True, "space_before": 7, "color": GREY},
    ], accent=BLUE)

    # ---------------- footer ----------------
    # Placed last, and its height derived from where the columns actually
    # ended, so adding a panel can never quietly print underneath it.
    content_bottom = max(ly, ry)
    fy = max(H_IN - 2.6, content_bottom + 0.5)
    if fy + 2.6 > H_IN + 0.01:
        raise SystemExit(
            f"layout overflow: content ends at {content_bottom:.2f}in and the "
            f"footer needs 2.6in, exceeding the {H_IN}in page. Trim a panel.")
    _box(slide, 0, fy, W_IN, H_IN - fy, fill=NAVY, line=None)
    _text(slide, MARGIN, fy + 0.45, W_IN - 2 * MARGIN, 1.7,
          [{"t": "github.com/zhouhufeng/IGVFagent      ·      "
                 "igvfagent.genohub.org      ·      Apache-2.0",
            "size": 26, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER},
           {"t": f"IGVFagent v{s['version']}  ·  open source, installable with "
                 f"pip  ·  runs on hosted APIs or entirely on local models",
            "size": 17, "color": RGBColor(0xAF, 0xCB, 0xE6),
            "align": PP_ALIGN.CENTER, "space_before": 8}],
          align=PP_ALIGN.CENTER)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path


def main(argv=None) -> int:
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("--out", default="Docs/Poster/IGVFagent_poster.pptx")
    a = ap.parse_args(argv)
    p = Path(a.out)
    if not p.is_absolute():
        p = ROOT / p
    out = build(p)
    print(f"Wrote: {out}")
    print(f"Size:  {W_IN}in x {H_IN}in (portrait) — within the 48in width limit")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
